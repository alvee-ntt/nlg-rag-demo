from __future__ import annotations

import csv
import io
from pathlib import Path

from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from .pdf_layout import Section, extract_sections, render_markdown

SUPPORTED_EXTENSIONS = {
    ".csv",
    ".docx",
    ".html",
    ".htm",
    ".md",
    ".pdf",
    ".pptx",
    ".txt",
    ".xlsx",
}


def is_supported(blob_name: str) -> bool:
    return Path(blob_name).suffix.lower() in SUPPORTED_EXTENSIONS


def extract_text(blob_name: str, content: bytes) -> str:
    ext = Path(blob_name).suffix.lower()
    if ext in {".txt", ".md"}:
        return _decode_text(content)
    if ext in {".html", ".htm"}:
        soup = BeautifulSoup(_decode_text(content), "html.parser")
        return soup.get_text("\n", strip=True)
    if ext == ".pdf":
        return _extract_pdf(content)
    if ext == ".docx":
        return _extract_docx(content)
    if ext == ".pptx":
        return _extract_pptx(content)
    if ext == ".xlsx":
        return _extract_xlsx(content)
    if ext == ".csv":
        return _extract_csv(content)
    raise ValueError(f"Unsupported file type: {ext}")


def extract_document(blob_name: str, content: bytes) -> list[Section]:
    """Extract a document as structured sections carrying page and heading provenance.

    PDFs go through layout reconstruction. Every other format is already in reading
    order, so it becomes a single unheaded section and is chunked the same way.
    """
    if Path(blob_name).suffix.lower() == ".pdf":
        return extract_sections(content)
    text = extract_text(blob_name, content)
    return [Section(text=text, page=1, heading_path=())] if text.strip() else []


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp1252"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="ignore")


def _extract_pdf(content: bytes) -> str:
    return render_markdown(extract_sections(content))


def _extract_docx(content: bytes) -> str:
    doc = Document(io.BytesIO(content))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    parts = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        if slide_parts:
            parts.append(f"[slide {slide_number}]\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def _extract_xlsx(content: bytes) -> str:
    workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts = []
    for sheet in workbook.worksheets:
        parts.append(f"[sheet {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                parts.append(" | ".join(values))
    return "\n".join(parts)


def _extract_csv(content: bytes) -> str:
    text = _decode_text(content)
    output = []
    for row in csv.reader(io.StringIO(text)):
        values = [cell.strip() for cell in row if cell.strip()]
        if values:
            output.append(" | ".join(values))
    return "\n".join(output)
